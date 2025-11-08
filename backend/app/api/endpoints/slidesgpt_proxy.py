from fastapi import APIRouter, Request, Form
from fastapi.responses import JSONResponse
import httpx
import logging
from app.core.config import SLIDESGPT_API_KEY
from fastapi.responses import StreamingResponse
from pptx import Presentation
from io import BytesIO
from pptx.util import Inches, Pt
import json
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


router = APIRouter()

@router.post("/slidesgpt/generate")
async def slidesgpt_generate(request: Request):
    data = await request.json()
    prompt = data.get("prompt")
    logging.info(f"prompt used is \n\n {prompt}")

    
    api_key = SLIDESGPT_API_KEY
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    payload = {"prompt": prompt}
    try:
        async with httpx.AsyncClient(timeout=130.0) as client:
            response = await client.post(
                "https://api.slidesgpt.com/v1/presentations/generate",
                headers=headers,
                json=payload
            )
            # Only read text if response is not OK
            if response.status_code >= 200 and response.status_code < 300:
                return JSONResponse(content=response.json(), status_code=response.status_code)
            else:
                error_str = await response.aread()
                error_str = error_str.decode(errors="replace")
                return JSONResponse(
                    content={"detail": "SlidesGPT API failed!", "error": error_str},
                    status_code=response.status_code
                )
    except httpx.ReadTimeout:
        return JSONResponse(
            content={"detail": "SlidesGPT API timed out! Try again later."},
            status_code=504
        )
    except Exception as e:
        return JSONResponse(
            content={"detail": f"SlidesGPT API error: {str(e)}"},
            status_code=500
        )
    
@router.get("/slidesgpt/download/{presentation_id}")
async def slidesgpt_download(presentation_id: str):
    api_key = SLIDESGPT_API_KEY
    url = f"https://api.slidesgpt.com/v1/presentations/{presentation_id}/download"
    headers = {"Authorization": f"Bearer {api_key}"}
    async with httpx.AsyncClient(timeout=None) as client:
        resp = await client.get(url, headers=headers)
        # Check for error
        if resp.status_code != 200:
            return JSONResponse(content={"detail": "Download failed!", "error": await resp.aread()}, status_code=resp.status_code)
        # Stream the pptx back
        return StreamingResponse(resp.aiter_bytes(), media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={
            "Content-Disposition": f"attachment; filename=presentation-{presentation_id}.pptx"
        })


@router.post("/slidesgpt/add-links")
async def add_links_to_pptx(
    presentation_id: str = Form(...),
    source_urls: str = Form(...),  # JSON-stringified list
):
    # Download PPTX
    api_key = SLIDESGPT_API_KEY
    url = f"https://api.slidesgpt.com/v1/presentations/{presentation_id}/download"
    headers = {"Authorization": f"Bearer {api_key}"}
    async with httpx.AsyncClient(timeout=None) as client:
        resp = await client.get(url, headers=headers)
        if resp.status_code != 200:
            return JSONResponse(content={"detail": "Download failed!", "error": await resp.aread()}, status_code=resp.status_code)
        pptx_bytes = await resp.aread()

    # Parse URLs and patch PPTX
    urls = json.loads(source_urls)
    prs = Presentation(BytesIO(pptx_bytes))
    slide_width = prs.slide_width      # in EMUs
    slide_height = prs.slide_height    # in EMUs

    # ----------- REMOVE "Photo by ..." CAPTIONS -----------
    for slide in prs.slides:
        to_remove = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip().lower()
            if "photo by" in text:
                to_remove.append(shape)
        for shape in to_remove:
            slide.shapes._spTree.remove(shape._element)

    # Convert 0.5 inches to EMUs for margin/padding
    left_margin = Inches(0.5)
    right_margin = Inches(0.5)
    footer_height = Inches(0.35)
    bottom_margin = Inches(0.3)   # space from bottom of slide

    # Placement (left, top, width, height) for textbox
    left = left_margin
    width = slide_width - left_margin - right_margin
    top = slide_height - bottom_margin - footer_height
    height = footer_height
    for idx, slide in enumerate(prs.slides):
        if idx < 3:
            continue  # Skip first 3 slides
        url_idx = idx - 3  # The first URL in your list goes to slide 4 (index 3)
        if url_idx < len(urls) and urls[url_idx]:
            url = urls[url_idx]
            #left = Inches(0.5)
            #top = Inches(6.0)
            #width = Inches(9)
            #height = Inches(0.4)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = url
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0, 102, 204)
            run.font.italic = True
            run.hyperlink.address = url

    out_bytes = BytesIO()
    prs.save(out_bytes)
    out_bytes.seek(0)
    return StreamingResponse(
        out_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f"attachment; filename=presentation-{presentation_id}-with-links.pptx"
        }
    )