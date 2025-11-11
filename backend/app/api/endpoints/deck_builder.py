# app/api/endpoints/deck_builder.py
from fastapi import APIRouter, Request, Query, Depends
from fastapi.responses import StreamingResponse, JSONResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from io import BytesIO
from datetime import datetime
from pathlib import Path
from urllib.parse import quote
import os
import shutil
import subprocess
import logging

# Set this to your public backend base when you have HTTPS; stays http://localhost:8000 in dev
from app.core.config import BACKEND_BASE
from app.dependencies import get_current_user

router = APIRouter()

# Template mapping
TEMPLATE_DIR = Path("templates")
TEMPLATE_MAP = {
    "System Default Template": None,  # Use vanilla PowerPoint (default)
    "Standard Research Template": "standard_research.pptx",
    "Executive Summary Template": "executive_summary.pptx",
    "Detailed Analysis Template": "detailed_analysis.pptx",
}

# Where generated decks are stored (main.py must mount: app.mount("/static", StaticFiles(directory="static"), name="static"))
STATIC_DECK_DIR = Path("static/decks")
STATIC_DECK_DIR.mkdir(parents=True, exist_ok=True)

logger = logging.getLogger(__name__)

# Helper function to format timestamp as DDMonYYYY_HHMM
def format_readable_timestamp(dt: datetime) -> str:
    """Format datetime as '10Nov2024_1430'"""
    months = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
    day = str(dt.day).zfill(2)
    month = months[dt.month - 1]
    year = dt.year
    hour = str(dt.hour).zfill(2)
    minute = str(dt.minute).zfill(2)
    return f"{day}{month}{year}_{hour}{minute}"

# -------------------- Helpers: layout & placeholders --------------------

def _find_layout_with_title(prs: Presentation):
    """Find a slide layout that has a real TITLE/CENTER_TITLE placeholder; else return layout[0]."""
    for layout in prs.slide_layouts:
        for shp in getattr(layout, "shapes", []):
            if getattr(shp, "is_placeholder", False):
                try:
                    if shp.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                        return layout
                except Exception:
                    pass
    return prs.slide_layouts[0]

def _find_layout_with_body(prs: Presentation):
    """Find a slide layout that has a BODY placeholder; else return layout[1] or layout[0]."""
    for layout in prs.slide_layouts:
        for shp in getattr(layout, "shapes", []):
            if getattr(shp, "is_placeholder", False):
                try:
                    if shp.placeholder_format.type == PP_PLACEHOLDER.BODY:
                        return layout
                except Exception:
                    pass
    return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

# NEW: Prefer a true Blank layout for content slides
def _find_blank_layout(prs: Presentation):
    # Try by name
    for layout in prs.slide_layouts:
        name = getattr(layout, "name", "") or ""
        if name.strip().lower() == "blank":
            return layout
    # Try by “no placeholders” heuristic
    for layout in prs.slide_layouts:
        try:
            if not any(getattr(s, "is_placeholder", False) for s in layout.shapes):
                return layout
        except Exception:
            pass
    return None  # caller will fall back

def _get_title_placeholder(slide):
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                return ph
        except Exception:
            continue
    return None

def _get_subtitle_placeholder(slide):
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.type in (PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY):
                return ph
        except Exception:
            continue
    return None

# NEW: remove placeholders that might overlap our custom boxes
def _remove_placeholders(slide, types=(PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY)):
    for shp in list(slide.shapes):
        if getattr(shp, "is_placeholder", False):
            try:
                if shp.placeholder_format.type in types:
                    el = shp._element
                    el.getparent().remove(el)
            except Exception:
                continue
# -------------------- Robust title writer (returns next Y) --------------------

def set_title_text(
    slide,
    text: str,
    *,
    left: float = 0.7,   # inches
    top: float = 0.4,    # inches
    width: float = 8.5,  # inches
    height: float = 1.2, # inches (nominal)
    force_box: bool = False,
) -> float:
    """
    Safe title writer with dynamic sizing & wrapping to avoid overlaps.
    - Uses slide.shapes.title when available (unless force_box=True).
    - Else draws a textbox at (left, top, width, height) in inches.
    - Returns a recommended 'next_top' (inches) to place the next element below the title.
    """
    title_shape = None if force_box else slide.shapes.title
    if title_shape is None:
        title_shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

    tf = title_shape.text_frame
    tf.word_wrap = True
    tf.clear()

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = (text or "").strip()

    # Dynamic font sizing that’s conservative to prevent overlap
    n = len(r.text)
    if n > 100:
        fsz = Pt(22); lines = 3
    elif n > 80:
        fsz = Pt(24); lines = 3
    elif n > 60:
        fsz = Pt(28); lines = 2
    elif n > 40:
        fsz = Pt(32); lines = 2
    else:
        fsz = Pt(36); lines = 1

    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = fsz
            run.font.bold = True

    # Estimate consumed height and return the next suggested top (inches)
    est_per_line = 0.55
    consumed = max(height, lines * est_per_line)
    next_top = top + consumed + 0.20  # + padding
    return next_top

# -------------------- Footer link --------------------

def add_footer_link(slide, prs, text, url):
    """Adds a footer at bottom with a clickable URL."""
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    left = Inches(0.5)
    right_margin = Inches(0.5)
    footer_height = Inches(0.35)
    bottom_margin = Inches(0.3)

    width = slide_w - left - right_margin
    top = slide_h - bottom_margin - footer_height
    height = footer_height

    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0, 102, 204)
    if url:
        run.hyperlink.address = url

# -------------------- Endpoint --------------------

@router.post("/deck/build-ppt")
async def build_ppt(
    request: Request,
    return_url: int = Query(0),   # if 1, save & return URLs instead of streaming
    user = Depends(get_current_user),  # Get current user for username
):
    """
    JSON body:
    {
      "title": "Gen AI & LLM Trends - July 2025",
      "subtitle": "Your optional subtitle",
      "template": "Standard Research Template",
      "include_summary": true,
      "sections": [
        {
          "category": "AI & GenAI Trends",
          "articles": [
            {
              "main_title": "Concise slide title",
              "original_title": "Original Article Title",
              "summary": "Full summary text...",
              "source": "https://example.com/path"
            }
          ]
        }
      ]
    }
    """
    try:
        data = await request.json()
    except Exception:
        return JSONResponse({"detail": "Invalid JSON"}, status_code=400)

    title = data.get("title") or "Research Deck"
    subtitle = data.get("subtitle") or datetime.utcnow().strftime("%d-%b-%Y")
    include_summary = bool(data.get("include_summary", True))
    sections = data.get("sections") or []
    template_name = data.get("template") or "System Default Template"

    # Load selected template
    template_file = TEMPLATE_MAP.get(template_name)

    if template_file is None:
        # System Default Template - use vanilla PowerPoint
        prs = Presentation()
        logger.info("Using System Default Template (vanilla PowerPoint)")
    elif template_file:
        template_path = TEMPLATE_DIR / template_file
        if template_path.exists():
            prs = Presentation(str(template_path))
            logger.info(f"Loaded template: {template_file}")

            # Clear all existing slides from template, keep only layouts/theme
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
            logger.info(f"Cleared template slides, using only theme/layouts")
        else:
            logger.warning(f"Template {template_file} not found; using default theme.")
            prs = Presentation()
    else:
        logger.warning(f"Unknown template '{template_name}'; using default theme.")
        prs = Presentation()

    # ----- Title slide -----
    title_layout = _find_layout_with_title(prs)
    s0 = prs.slides.add_slide(title_layout)

    # Title text (robust)
    title_ph = _get_title_placeholder(s0)
    if title_ph:
        # use the placeholder, but still run through sizing helper to get next_top
        next_top = set_title_text(s0, title)  # uses placeholder by default
    else:
        # draw our own box at a consistent spot
        next_top = set_title_text(s0, title, force_box=True, left=0.7, top=0.4)

    # Subtitle text (robust)
    subph = _get_subtitle_placeholder(s0)
    if subph and getattr(subph, "has_text_frame", False):
        subph.text = subtitle
    else:
        # place a textbox just below the title
        subtitle_top = max(next_top, 1.4)
        tb2 = s0.shapes.add_textbox(Inches(0.7), Inches(subtitle_top), Inches(8.5), Inches(0.8))
        tf2 = tb2.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(20)

    # ----- Executive summary (optional) -----
    if include_summary:
        bullet_layout = _find_layout_with_body(prs)
        s_sum = prs.slides.add_slide(bullet_layout)

        # Title for summary slide; always use our own box to avoid weird templates
        sum_title_bottom = set_title_text(s_sum, "Executive Summary", force_box=True, left=0.7, top=0.4)

        # Body: use BODY placeholder if present, else textbox
        body_ph = None
        for ph in s_sum.placeholders:
            try:
                if ph.placeholder_format.type == PP_PLACEHOLDER.BODY:
                    body_ph = ph
                    break
            except Exception:
                pass

        total_articles = sum(len(sec.get("articles", [])) for sec in sections)
        if body_ph and body_ph.has_text_frame:
            tf = body_ph.text_frame
            tf.word_wrap = True
            tf.clear()
            p1 = tf.paragraphs[0]
            p1.text = f"Total articles: {total_articles}"
            for sec in sections:
                p = tf.add_paragraph()
                p.text = f"{sec.get('category')}: {len(sec.get('articles', []))}"
                p.level = 1
        else:
            body_top = max(sum_title_bottom, 1.6)
            tx = s_sum.shapes.add_textbox(Inches(0.7), Inches(body_top), Inches(8.0), Inches(4.6))
            tf = tx.text_frame
            tf.word_wrap = True
            tf.clear()
            p1 = tf.paragraphs[0]
            p1.text = f"Total articles: {total_articles}"
            for sec in sections:
                p = tf.add_paragraph()
                p.text = f"{sec.get('category')}: {len(sec.get('articles', []))}"
                p.level = 1

     # Content slides — use Blank if possible, then purge placeholders
    blank_layout = _find_blank_layout(prs)
    content_layout = blank_layout or _find_layout_with_body(prs)

    for sec in sections:
        for art in sec.get("articles", []):
            slide = prs.slides.add_slide(content_layout)
            # Safety: remove any built-in placeholders that would overlap
            _remove_placeholders(slide)

            # 1) Title — always use our own box; templates vary wildly
            main_title = (art.get("main_title") or art.get("original_title") or "Untitled")[:200]
            title_bottom = set_title_text(slide, main_title, force_box=True, left=0.7, top=0.6)

            # 2) Subtitle — just below the title
            subtitle_top = max(title_bottom, 1.6)
            subtitle_box = slide.shapes.add_textbox(Inches(0.7), Inches(subtitle_top), Inches(8.0), Inches(0.8))
            stf = subtitle_box.text_frame
            stf.word_wrap = True
            stf.clear()
            p_sub = stf.paragraphs[0]
            p_sub.text = (art.get("original_title") or "").strip()
            for r in p_sub.runs:
                r.font.size = Pt(16)

            # 3) Body — under the subtitle
            body_top = subtitle_top + 0.9
            tx = slide.shapes.add_textbox(Inches(0.7), Inches(body_top), Inches(8.0), Inches(4.2))
            tf = tx.text_frame
            tf.word_wrap = True
            tf.clear()
            summary = (art.get("summary") or "").strip()
            for i, block in enumerate(summary.split("\n")):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = block.strip()
                p.level = 0

            # 4) Footer source
            src = art.get("source") or ""
            if src:
                visible = src.replace("https://", "").replace("http://", "").replace("www.", "")
                add_footer_link(slide, prs, f"Source: {visible}", src)

    # ----- Output handling -----
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)

    # Generate filename with username and readable timestamp
    # Format: Title_username_10Nov2024_1430.pptx
    now = datetime.utcnow()
    readable_ts = format_readable_timestamp(now)
    username = user.username if user else "unknown"
    base_filename = f"{title.replace(' ', '_')}_{username}_{readable_ts}.pptx"

    if not return_url:
        # Stream the PPTX back for download
        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{base_filename}"'}
        )

    # Save PPTX to disk for preview
    disk_pptx = STATIC_DECK_DIR / base_filename
    with open(disk_pptx, "wb") as f:
        f.write(buf.getbuffer())

    file_url = f"/static/decks/{quote(base_filename)}"
    absolute_file_url = f"{BACKEND_BASE}{file_url}"

    # Try PDF conversion (LibreOffice) for local embed fallback
    pdf_abs_url = None
    pdf_rel_url = None
    try:
        soffice = shutil.which("soffice")
        if not soffice:
            mac_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
            soffice = mac_path if os.path.exists(mac_path) else None

        if soffice:
            cmd = [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(STATIC_DECK_DIR), str(disk_pptx)]
            logger.info("Converting PPTX to PDF with: %s", " ".join(cmd))
            subprocess.run(cmd, check=True, timeout=90)
            filename_pdf = base_filename.replace(".pptx", ".pdf")
            pdf_rel_url = f"/static/decks/{quote(filename_pdf)}"
            pdf_abs_url = f"{BACKEND_BASE}{pdf_rel_url}"
        else:
            logger.warning("LibreOffice not found; PDF preview disabled.")
    except Exception as e:
        logger.warning("PDF conversion failed: %s", e)

    # Office Online viewer only works with a public HTTPS absolute URL
    viewer_url = (
        f"https://view.officeapps.live.com/op/embed.aspx?src={quote(absolute_file_url, safe='')}"
        if absolute_file_url.startswith("https://")
        else None
    )

    return JSONResponse(
        {
            "file_url": file_url,                # relative PPTX (download)
            "absolute_file_url": absolute_file_url,
            "viewer_url": viewer_url,            # Office Online (if https)
            "pdf_url": pdf_rel_url,              # relative PDF (if converted)
            "pdf_abs_url": pdf_abs_url           # absolute PDF (good for local iframe)
        },
        status_code=200
    )